from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 흰 배경 + #3A59D1 테마 팔레트 ───────────────
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # 흰색 (헤더 텍스트용)
DARK_BG   = RGBColor(0xFF, 0xFF, 0xFF)   # 슬라이드 배경 = 흰색
DARK_TEXT = RGBColor(0x1A, 0x2B, 0x55)   # 본문 텍스트 (진한 네이비)
ACCENT    = RGBColor(0x3A, 0x59, 0xD1)   # 주요 강조색 #3A59D1
GREEN     = RGBColor(0x1A, 0x7A, 0x4A)   # 진한 초록
YELLOW    = RGBColor(0xB8, 0x6C, 0x00)   # 진한 앰버
RED       = RGBColor(0xC0, 0x39, 0x2B)   # 진한 빨강
SUBTEXT   = RGBColor(0x5A, 0x6A, 0x9E)   # 중간 블루-그레이
CARD_BG   = RGBColor(0xEE, 0xF2, 0xFF)   # 연한 블루 카드
HEADER_BG = RGBColor(0x3A, 0x59, 0xD1)   # 헤더 배경 #3A59D1
ROW_ALT   = RGBColor(0xE0, 0xE8, 0xFF)   # 교차 행 연블루
CARD_HDR  = RGBColor(0x2A, 0x45, 0xB8)   # 카드 헤더 (약간 어두운 블루)

blank_layout = prs.slide_layouts[6]


def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=None,
                align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = DARK_TEXT
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, HEADER_BG)
    add_textbox(slide, title, 0.4, 0.1, 10, 0.6, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.4, 0.65, 12, 0.35,
                    font_size=13, color=RGBColor(0xCC, 0xDA, 0xFF))


# ──────────────────────────────────────────────
# Slide 1: Title
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

add_rect(slide, 1.5, 1.8, 10.3, 0.08, ACCENT)
add_textbox(slide, "메인 에이전트 (통합 분석)", 1.5, 2.1, 10.3, 1.2,
            font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(slide,
            "뉴스 · 공시 · 재무 서브에이전트를 병렬 호출하여\n종합 투자 시그널과 LLM 분석 요약을 반환하는 오케스트레이터 에이전트",
            1.5, 3.5, 10.3, 1.0, font_size=16, color=SUBTEXT, align=PP_ALIGN.CENTER)

labels = [("gpt-5-mini", GREEN), ("LangGraph RAG", ACCENT), ("PostgreSQL 캐시", YELLOW)]
x = 3.2
for label, col in labels:
    add_rect(slide, x, 5.0, 2.2, 0.45, CARD_BG)
    add_textbox(slide, label, x, 5.0, 2.2, 0.45,
                font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    x += 2.5

add_textbox(slide, "antelligen-backend  |  2026", 0.4, 6.9, 12.5, 0.4,
            font_size=11, color=SUBTEXT, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────
# Slide 2: 아키텍처 흐름
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_header(slide, "아키텍처 흐름", "POST /api/v1/agent/query 요청 처리 과정")

steps = [
    ("① 인증 확인", "쿠키 / 쿼리 파라미터 / Authorization 헤더", ACCENT),
    ("② 캐시 조회", "PostgreSQL — 동일 ticker 1시간 이내 결과 재사용", YELLOW),
    ("③ 병렬 호출", "asyncio.gather() → 뉴스 · 공시 · 재무 동시 실행", GREEN),
    ("④ 시그널 집계", "가중 평균: Σ(signal × confidence) / Σ(confidence)", ACCENT),
    ("⑤ LLM 종합", "OpenAISynthesisClient (gpt-5-mini) → summary + key_points", YELLOW),
    ("⑥ 저장 및 반환", "integrated_analysis_results INSERT → FrontendAgentResponse", GREEN),
]

for i, (title, desc, col) in enumerate(steps):
    row = i // 2
    col_idx = i % 2
    lx = 0.4 + col_idx * 6.5
    ty = 1.4 + row * 1.85

    add_rect(slide, lx, ty, 6.1, 1.6, CARD_BG)
    add_rect(slide, lx, ty, 0.07, 1.6, col)
    add_textbox(slide, title, lx + 0.2, ty + 0.12, 5.8, 0.45,
                font_size=15, bold=True, color=col)
    add_textbox(slide, desc, lx + 0.2, ty + 0.65, 5.8, 0.75,
                font_size=12, color=SUBTEXT)


# ──────────────────────────────────────────────
# Slide 3: 서브에이전트 상세
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_header(slide, "서브에이전트 상세", "각 에이전트의 데이터 소스 · 기술 · 응답 시간")

agents = [
    ("뉴스 에이전트", GREEN,
     "• 네이버 뉴스 API 수집 (8개 종목)\n• OpenAI gpt-5-mini 감성 분석\n• 최신 기사 우선 가중\n• DB 없으면 자동 수집 후 재분석",
     "8 ~ 12 초"),
    ("공시 에이전트", YELLOW,
     "• DART API → 공시 수집 및 저장\n• LangGraph RAG (gpt-5-mini)\n• 3가지 분석 타입 (full/signal/flow)\n• Redis 캐시 1시간",
     "캐시 ~500ms / RAG 7~20초"),
    ("재무 에이전트", ACCENT,
     "• SerpAPI + DART 재무비율 수집\n• pgvector 코사인 유사도 검색\n• LangGraph RAG (gpt-5-mini)\n• 첫 종목 자동 수집",
     "캐시 ~1초 / 첫 분석 40~60초"),
]

for i, (name, col, desc, timing) in enumerate(agents):
    lx = 0.4 + i * 4.3
    add_rect(slide, lx, 1.3, 4.0, 5.7, CARD_BG)
    add_rect(slide, lx, 1.3, 4.0, 0.5, col)
    add_textbox(slide, name, lx, 1.3, 4.0, 0.5,
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, lx + 0.15, 2.0, 3.75, 3.2,
                font_size=12, color=DARK_TEXT)
    add_rect(slide, lx + 0.2, 5.45, 3.6, 0.06, col)
    add_textbox(slide, "⏱  " + timing, lx + 0.15, 5.6, 3.75, 0.9,
                font_size=12, color=col)


# ──────────────────────────────────────────────
# Slide 4: 엔드포인트 스펙
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_header(slide, "엔드포인트 스펙")

add_rect(slide, 0.4, 1.3, 5.9, 5.7, CARD_BG)
add_textbox(slide, "POST  /api/v1/agent/query", 0.5, 1.38, 5.7, 0.45,
            font_size=14, bold=True, color=GREEN)
req = (
    "요청 Body:\n"
    "  ticker     종목 코드 (예: 005930)\n"
    "  query      사용자 질문\n\n"
    "응답 주요 필드:\n"
    "  result_status  success | partial_failure | failure\n"
    "  answer         LLM 종합 요약\n"
    "  agent_results  서브에이전트별 결과 배열\n"
    "    ├ agent_name   news | disclosure | finance\n"
    "    ├ signal       bullish | bearish | neutral\n"
    "    ├ confidence   0.0 ~ 1.0\n"
    "    ├ summary      에이전트 분석 요약\n"
    "    └ key_points   핵심 포인트 리스트"
)
add_textbox(slide, req, 0.5, 1.92, 5.7, 4.8, font_size=11.5, color=SUBTEXT)

add_rect(slide, 7.0, 1.3, 5.9, 2.7, CARD_BG)
add_textbox(slide, "GET  /api/v1/agent/history", 7.1, 1.38, 5.7, 0.45,
            font_size=14, bold=True, color=ACCENT)
hist = (
    "파라미터:\n"
    "  ticker  종목 코드 (필수)\n"
    "  limit   최대 50건 (기본 10)\n\n"
    "응답: 최근 분석 이력 배열\n"
    "  overall_signal, confidence,\n"
    "  summary, key_points, created_at"
)
add_textbox(slide, hist, 7.1, 1.92, 5.7, 2.2, font_size=11.5, color=SUBTEXT)

add_rect(slide, 7.0, 4.15, 5.9, 2.85, CARD_BG)
add_textbox(slide, "인증 방식 (공통)", 7.1, 4.22, 5.7, 0.45,
            font_size=14, bold=True, color=YELLOW)
auth = (
    "① 쿼리 파라미터:  ?token={uuid}\n"
    "② 쿠키:           user_token\n"
    "③ 헤더:           Authorization: Bearer {uuid}\n\n"
    "미인증 시 401 → 프론트 로그인 페이지 redirect"
)
add_textbox(slide, auth, 7.1, 4.78, 5.7, 2.1, font_size=11.5, color=SUBTEXT)


# ──────────────────────────────────────────────
# Slide 5: 시그널 집계 & DB 테이블
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_header(slide, "시그널 집계 & DB 테이블")

add_rect(slide, 0.4, 1.3, 5.9, 5.7, CARD_BG)
add_textbox(slide, "시그널 가중 집계 로직", 0.5, 1.38, 5.7, 0.45,
            font_size=14, bold=True, color=ACCENT)
add_textbox(slide,
            "score = Σ(signal_score × confidence)\n         / Σ(confidence)",
            0.5, 2.0, 5.7, 0.9, font_size=13, bold=True, color=DARK_TEXT)

for i, (val, label, col) in enumerate([
    ("+1.0", "bullish  →  매수 신호", GREEN),
    (" 0.0", "neutral  →  관망 신호", SUBTEXT),
    ("-1.0", "bearish  →  매도 신호", RED),
]):
    ty = 3.1 + i * 0.6
    add_rect(slide, 0.5, ty, 5.6, 0.5, ROW_ALT)
    add_textbox(slide, val, 0.7, ty + 0.04, 1.0, 0.42,
                font_size=13, bold=True, color=col)
    add_textbox(slide, label, 1.8, ty + 0.04, 4.0, 0.42,
                font_size=13, color=DARK_TEXT)

thresholds = (
    "판정 기준:\n"
    "  score > +0.2  →  bullish\n"
    "  score < −0.2  →  bearish\n"
    "  그 외          →  neutral"
)
add_textbox(slide, thresholds, 0.5, 4.95, 5.7, 1.8, font_size=12, color=SUBTEXT)

add_rect(slide, 7.0, 1.3, 5.9, 5.7, CARD_BG)
add_textbox(slide, "integrated_analysis_results", 7.1, 1.38, 5.7, 0.45,
            font_size=14, bold=True, color=ACCENT)

cols_data = [
    ("ticker",         "VARCHAR",  "종목 코드 (인덱스)"),
    ("overall_signal", "VARCHAR",  "bullish / bearish / neutral"),
    ("confidence",     "FLOAT",    "0.0 ~ 1.0"),
    ("summary",        "TEXT",     "LLM 종합 요약"),
    ("key_points",     "JSON",     "핵심 포인트 리스트"),
    ("sub_results",    "JSON",     "서브에이전트 결과 전체"),
    ("created_at",     "DATETIME", "생성 시각 (캐시 기준)"),
]
for i, (col, typ, desc) in enumerate(cols_data):
    ty = 1.95 + i * 0.62
    bg = CARD_BG if i % 2 == 0 else ROW_ALT
    add_rect(slide, 7.05, ty, 5.8, 0.56, bg)
    add_textbox(slide, col,  7.15, ty + 0.06, 1.8, 0.44,
                font_size=11, bold=True, color=ACCENT)
    add_textbox(slide, typ,  9.0,  ty + 0.06, 1.2, 0.44,
                font_size=10, color=YELLOW)
    add_textbox(slide, desc, 10.2, ty + 0.06, 2.5, 0.44,
                font_size=10, color=SUBTEXT)

add_textbox(slide, "※ 캐시: created_at 기준 1시간 이내 재사용",
            7.1, 6.45, 5.7, 0.4, font_size=11, color=SUBTEXT)


# ──────────────────────────────────────────────
# Slide 6: 에러 처리 방식
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_header(slide, "에러 처리 방식", "서브에이전트 실패 격리 · result_status 판정 · 집계 제외 규칙")

# 왼쪽: 실패 격리 흐름
add_rect(slide, 0.35, 1.3, 5.9, 5.7, CARD_BG)
add_rect(slide, 0.35, 1.3, 5.9, 0.42, RED)
add_textbox(slide, "서브에이전트 실패 격리", 0.35, 1.3, 5.9, 0.42,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, "asyncio.gather(return_exceptions=True)",
            0.5, 1.88, 5.6, 0.38, font_size=11, bold=True, color=ACCENT)
add_textbox(slide, "하나가 실패해도 나머지 에이전트는 계속 실행",
            0.5, 2.22, 5.6, 0.3, font_size=10, color=SUBTEXT)

agents_err = [
    ("뉴스",  "DB 없음 → 자동 수집 후 재시도\n재시도 후에도 없으면 no_data 반환",  GREEN),
    ("공시",  "예외 발생 시 error 반환\nRedis 캐시 히트 시 즉시 반환",             YELLOW),
    ("재무",  "벡터 DB 없음 → 자동 수집 후 재시도\n수집 실패 시 error 반환",       ACCENT),
]
for i, (name, desc, col) in enumerate(agents_err):
    ty = 2.68 + i * 1.08
    add_rect(slide, 0.45, ty, 5.7, 0.95, ROW_ALT if i % 2 else CARD_BG)
    add_rect(slide, 0.45, ty, 0.06, 0.95, col)
    add_textbox(slide, name + " 에이전트", 0.6, ty + 0.06, 1.8, 0.35,
                font_size=11, bold=True, color=col)
    add_textbox(slide, desc, 0.6, ty + 0.44, 5.4, 0.44,
                font_size=10, color=SUBTEXT)

add_textbox(slide, "실패 에이전트 → _coerce() → SubAgentResponse.error()",
            0.5, 5.98, 5.6, 0.38, font_size=10, color=SUBTEXT)

# 오른쪽 상단: result_status 판정
add_rect(slide, 6.6, 1.3, 6.38, 2.7, CARD_BG)
add_rect(slide, 6.6, 1.3, 6.38, 0.42, ACCENT)
add_textbox(slide, "result_status 판정", 6.6, 1.3, 6.38, 0.42,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

statuses = [
    ("3개 모두 성공",  "success",          GREEN),
    ("1~2개 성공",    "partial_failure",   YELLOW),
    ("전체 실패",     "failure",           RED),
]
for i, (cond, status, col) in enumerate(statuses):
    ty = 1.88 + i * 0.62
    bg = ROW_ALT if i % 2 else CARD_BG
    add_rect(slide, 6.7, ty, 6.15, 0.55, bg)
    add_textbox(slide, cond,   6.8,  ty + 0.1, 2.8, 0.35, font_size=11, color=DARK_TEXT)
    add_textbox(slide, status, 9.65, ty + 0.1, 2.9, 0.35, font_size=11, bold=True, color=col)

# 오른쪽 하단: 집계 & 저장 규칙
add_rect(slide, 6.6, 4.18, 6.38, 2.82, CARD_BG)
add_rect(slide, 6.6, 4.18, 6.38, 0.42, YELLOW)
add_textbox(slide, "집계 & 저장 규칙", 6.6, 4.18, 6.38, 0.42,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rules = [
    ("시그널 집계 제외",  "is_success() AND signal != None 조건 미충족 시\n실패/no_data 에이전트는 집계에서 제외"),
    ("DB 저장 조건",      "result_status == SUCCESS 일 때만 저장\npartial_failure / failure 는 저장 안 함"),
    ("캐시 재사용",       "저장된 결과는 1시간 동안 캐시\n동일 ticker 재요청 시 서브에이전트 미호출"),
]
for i, (title, desc) in enumerate(rules):
    ty = 4.75 + i * 0.72
    bg = ROW_ALT if i % 2 else CARD_BG
    add_rect(slide, 6.7, ty, 6.15, 0.65, bg)
    add_textbox(slide, title, 6.8, ty + 0.04, 2.2, 0.3, font_size=11, bold=True, color=ACCENT)
    add_textbox(slide, desc,  6.8, ty + 0.32, 5.9, 0.28, font_size=9,  color=SUBTEXT)


# ──────────────────────────────────────────────
# Slide 7: 분석 가능 종목 8개
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
add_slide_header(slide, "분석 가능 종목 (8개)",
                 "뉴스 에이전트 지원 종목 — 공시·재무 에이전트는 모든 종목 지원")

tickers = [
    ("005930", "삼성전자",        "반도체 · 전자",         '"삼성전자 투자해도 될까요?"'),
    ("000660", "SK하이닉스",      "HBM · 메모리",          '"HBM 실적 기대되는데 투자해도 될까?"'),
    ("005380", "현대차",          "전기차 · 모빌리티",     '"전기차 전환 어떻게 보고 있어?"'),
    ("035420", "네이버",          "AI · 플랫폼",           '"AI 사업 성장 가능성 어때?"'),
    ("035720", "카카오",          "플랫폼 · 콘텐츠",       '"지금 저점 매수 타이밍일까?"'),
    ("068270", "셀트리온",        "바이오시밀러",          '"바이오시밀러 전망은?"'),
    ("207940", "삼성바이오로직스","CMO · 바이오",          '"장기 투자 괜찮을까?"'),
    ("005490", "포스코",          "2차전지 · 철강",        '"2차전지 소재 사업 어때?"'),
]

card_colors = [GREEN, ACCENT, YELLOW, GREEN, ACCENT, YELLOW, GREEN, ACCENT]

for i, (ticker, name, sector, q) in enumerate(tickers):
    row = i // 4
    col_idx = i % 4
    lx = 0.3 + col_idx * 3.25
    ty = 1.3 + row * 2.9

    col = card_colors[i]
    add_rect(slide, lx, ty, 3.05, 2.55, CARD_BG)
    add_rect(slide, lx, ty, 3.05, 0.42, col)
    add_textbox(slide, name, lx, ty, 3.05, 0.42,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, ticker, lx + 0.1, ty + 0.52, 2.85, 0.35,
                font_size=12, bold=True, color=col)
    add_textbox(slide, sector, lx + 0.1, ty + 0.88, 2.85, 0.35,
                font_size=11, color=SUBTEXT)
    add_textbox(slide, q, lx + 0.1, ty + 1.35, 2.85, 0.95,
                font_size=10, color=DARK_TEXT)

add_textbox(slide,
            "※ 새 종목 추가: TICKER_TO_KEYWORDS & COLLECTION_KEYWORDS 두 곳 모두 수정",
            0.4, 7.05, 12.5, 0.35, font_size=11, color=SUBTEXT)


# ──────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────
output = "docs/main-agent-slides.pptx"
prs.save(output)
print(f"Saved: {output}")
