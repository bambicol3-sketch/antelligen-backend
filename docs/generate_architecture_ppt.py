from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x2B, 0x55)   # 딥 블루 배경
HEADER_BG = RGBColor(0x3A, 0x59, 0xD1)   # 헤더 블루
ACCENT    = RGBColor(0x5B, 0xC4, 0xFF)   # 밝은 블루
GREEN     = RGBColor(0xA6, 0xE3, 0xA1)   # green
YELLOW    = RGBColor(0xF9, 0xE2, 0xAF)   # yellow
RED       = RGBColor(0xF3, 0x8B, 0xA8)   # red/pink
PURPLE    = RGBColor(0xCB, 0xA6, 0xF7)   # purple
TEAL      = RGBColor(0x94, 0xE2, 0xD8)   # teal
SUBTEXT   = RGBColor(0xB8, 0xD4, 0xF0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG   = RGBColor(0x16, 0x3E, 0x75)   # 카드 블루
CARD_BG2  = RGBColor(0x1A, 0x48, 0x85)   # 카드 블루2
BORDER    = RGBColor(0x2D, 0x6A, 0xB0)   # 테두리 블루
DARK_TEXT = RGBColor(0x0D, 0x2B, 0x55)

blank_layout = prs.slide_layouts[6]


def set_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, t, w, h, fill, line_color=None, line_width=None):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt2(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, t, w, h,
        size=11, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(t), Inches(w), Inches(h)
    )
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def arrow(slide, left, t, w, h, color=SUBTEXT, vertical=True):
    """Draw a simple arrow connector (rectangle + triangle head)."""
    if vertical:
        rect(slide, left, t, w, h, color)
        # arrowhead triangle via thin rect
        rect(slide, left - w, t + h - 0.01, w * 3, 0.01, color)
    else:
        rect(slide, left, t, w, h, color)


# ══════════════════════════════════════════════
# Slide: System Architecture
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

# ── Header ───────────────────────────────────
rect(slide, 0, 0, 13.33, 0.75, HEADER_BG)
rect(slide, 0, 0.75, 13.33, 0.04, ACCENT)
txt(slide, "시스템 아키텍처", 0.35, 0.08, 7, 0.55,
    size=26, bold=True, color=WHITE)
txt(slide, "antelligen-backend  |  FastAPI + PostgreSQL + Redis + OpenAI",
    7.5, 0.18, 5.5, 0.4, size=11, color=SUBTEXT, align=PP_ALIGN.RIGHT)

# ── Layer labels (left side) ─────────────────
layer_info = [
    (0.82, "Client",       ACCENT),
    (1.72, "Frontend",     GREEN),
    (2.62, "Backend API",  YELLOW),
    (3.95, "AI Agents",    PURPLE),
    (5.6,  "Storage",      TEAL),
    (6.5,  "External API", RED),
]
for ty, label, col in layer_info:
    txt(slide, label, 0.05, ty, 0.9, 0.38,
        size=9, bold=True, color=col, align=PP_ALIGN.RIGHT)

# ── 가로 구분선 (레이어 경계) ─────────────────
dividers = [1.62, 2.52, 3.38, 5.42, 6.42]
for dy in dividers:
    rect(slide, 1.0, dy, 12.2, 0.015, BORDER)

# ══════════════════════════════════════════════
# ROW 1 — Client (y=0.85)
# ══════════════════════════════════════════════
rect(slide, 5.4, 0.85, 2.5, 0.62, CARD_BG2, ACCENT, 1.2)
txt(slide, "Browser / App", 5.4, 0.85, 2.5, 0.62,
    size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# ROW 2 — Frontend (y=1.75)
# ══════════════════════════════════════════════
rect(slide, 5.4, 1.75, 2.5, 0.62, CARD_BG2, GREEN, 1.2)
txt(slide, "React Frontend", 5.4, 1.75, 2.5, 0.62,
    size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# 화살표: client → frontend
rect(slide, 6.6, 1.49, 0.06, 0.26, SUBTEXT)
txt(slide, "▼", 6.54, 1.6, 0.2, 0.2, size=8, color=SUBTEXT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# ROW 3 — Backend API routers (y=2.65)
# ══════════════════════════════════════════════
routers = [
    (1.05, "Kakao Auth\n/kakao-auth",        YELLOW),
    (3.3,  "Agent Router\n/agent",           YELLOW),
    (5.55, "Auth Router\n/authentication",   YELLOW),
    (7.8,  "Board Router\n/board",           YELLOW),
    (10.05,"Stock Router\n/stock",           YELLOW),
]
for lx, label, col in routers:
    rect(slide, lx, 2.65, 2.0, 0.62, CARD_BG2, col, 0.8)
    txt(slide, label, lx, 2.65, 2.0, 0.62,
        size=10, bold=False, color=col, align=PP_ALIGN.CENTER)

# 화살표: frontend → backend
rect(slide, 6.6, 2.39, 0.06, 0.26, SUBTEXT)
txt(slide, "▼", 6.54, 2.5, 0.2, 0.2, size=8, color=SUBTEXT, align=PP_ALIGN.CENTER)
txt(slide, "HTTP / REST", 6.72, 2.43, 1.5, 0.2, size=8, color=SUBTEXT)

# ══════════════════════════════════════════════
# ROW 4 — AI Agents (y=3.52 ~ 5.3)
# ══════════════════════════════════════════════
# Main agent box
rect(slide, 1.05, 3.52, 11.2, 1.75, RGBColor(0x22, 0x24, 0x38), PURPLE, 1.0)
txt(slide, "ProcessAgentQueryUseCase  (메인 오케스트레이터)",
    1.15, 3.57, 11.0, 0.38, size=11, bold=True, color=PURPLE)

# Sub-agents inside
sub_agents = [
    (1.2,  "뉴스 에이전트\nNewsSubAgentAdapter\ngpt-5-mini",          GREEN),
    (4.85, "공시 에이전트\nDisclosureSubAgentAdapter\nLangGraph RAG",   YELLOW),
    (8.5,  "재무 에이전트\nFinanceSubAgentAdapter\npgvector + RAG",    ACCENT),
]
for lx, label, col in sub_agents:
    rect(slide, lx, 3.98, 3.45, 1.1, CARD_BG, col, 0.8)
    txt(slide, label, lx, 3.98, 3.45, 1.1,
        size=10, color=col, align=PP_ALIGN.CENTER)

# OpenAI synthesis
rect(slide, 5.5, 3.57, 2.3, 0.32, CARD_BG)
txt(slide, "OpenAISynthesisClient", 5.5, 3.57, 2.3, 0.32,
    size=9, color=PURPLE, align=PP_ALIGN.CENTER)

# asyncio.gather label
txt(slide, "asyncio.gather() ── 병렬 실행", 1.15, 3.85, 6.0, 0.28,
    size=9, color=SUBTEXT)

# 화살표: backend → main agent
rect(slide, 3.3, 3.28, 0.06, 0.24, SUBTEXT)
txt(slide, "▼", 3.24, 3.38, 0.2, 0.2, size=8, color=SUBTEXT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# ROW 5 — Storage (y=5.55)
# ══════════════════════════════════════════════
stores = [
    (1.05,  "PostgreSQL\n+ pgvector",          TEAL,   "통합분석 이력 / 공시 / 주식 벡터"),
    (4.7,   "Redis",                           TEAL,   "세션 토큰 / 공시 캐시 / 임시 토큰"),
    (8.35,  "PostgreSQL\n(뉴스 DB)",            TEAL,   "collected_news 테이블"),
]
for lx, label, col, sub in stores:
    rect(slide, lx, 5.55, 3.3, 0.75, CARD_BG2, col, 0.8)
    txt(slide, label, lx, 5.55, 3.3, 0.38,
        size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(slide, sub, lx + 0.1, 5.93, 3.1, 0.3,
        size=9, color=SUBTEXT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# ROW 6 — External APIs (y=6.48)
# ══════════════════════════════════════════════
ext_apis = [
    (1.05,  "Kakao OAuth",    RED),
    (3.25,  "Naver News API", RED),
    (5.45,  "DART API",       RED),
    (7.65,  "SerpAPI",        RED),
    (9.85,  "OpenAI API",     RED),
]
for lx, label, col in ext_apis:
    rect(slide, lx, 6.48, 1.95, 0.52, CARD_BG2, col, 0.7)
    txt(slide, label, lx, 6.48, 1.95, 0.52,
        size=10, bold=True, color=col, align=PP_ALIGN.CENTER)


# ── 저장 ────────────────────────────────────
output = "docs/system-architecture-slide.pptx"
prs.save(output)
print(f"Saved: {output}")
