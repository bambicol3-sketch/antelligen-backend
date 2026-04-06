from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 흰 배경 + #3A59D1 테마 팔레트 ───────────────
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG   = RGBColor(0xFF, 0xFF, 0xFF)   # 슬라이드 배경 = 흰색
DARK_TEXT = RGBColor(0x1A, 0x2B, 0x55)
ACCENT    = RGBColor(0x3A, 0x59, 0xD1)   # #3A59D1
GREEN     = RGBColor(0x1A, 0x7A, 0x4A)
YELLOW    = RGBColor(0xB8, 0x6C, 0x00)
RED       = RGBColor(0xC0, 0x39, 0x2B)
SUBTEXT   = RGBColor(0x5A, 0x6A, 0x9E)
CARD_BG   = RGBColor(0xEE, 0xF2, 0xFF)
HEADER_BG = RGBColor(0x3A, 0x59, 0xD1)
ROW_ALT   = RGBColor(0xE0, 0xE8, 0xFF)
LIGHT_HDR = RGBColor(0xCC, 0xDA, 0xFF)
PURPLE    = RGBColor(0x6A, 0x3A, 0xB8)
TEAL      = RGBColor(0x0A, 0x7A, 0x7A)

blank_layout = prs.slide_layouts[6]


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def tb(slide, text, l, t, w, h,
       size=12, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = DARK_TEXT
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rc(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def header(slide, title, subtitle=None):
    rc(slide, 0, 0, 13.33, 1.1, HEADER_BG)
    tb(slide, title, 0.4, 0.1, 10, 0.6, size=28, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, 0.4, 0.65, 12, 0.35, size=13, color=LIGHT_HDR)


# ══════════════════════════════════════════════
# Slide 1: Title
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)

rc(slide, 1.5, 1.8, 10.3, 0.08, ACCENT)
tb(slide, "UI 핵심 기능 소개", 1.5, 2.1, 10.3, 1.2,
   size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(slide, "stock-supporters-frontend  |  Next.js 16 · React 19 · Jotai · Tailwind CSS",
   1.5, 3.5, 10.3, 0.5, size=15, color=SUBTEXT, align=PP_ALIGN.CENTER)

badges = [
    ("Next.js 16", ACCENT), ("React 19", GREEN),
    ("Jotai", YELLOW), ("Tailwind CSS", TEAL), ("TypeScript", PURPLE),
]
x = 1.3
for label, col in badges:
    rc(slide, x, 4.8, 2.0, 0.45, CARD_BG)
    tb(slide, label, x, 4.8, 2.0, 0.45, size=12, bold=True,
       color=col, align=PP_ALIGN.CENTER)
    x += 2.2

tb(slide, "stock-supporters-frontend  |  2026", 0.4, 6.9, 12.5, 0.4,
   size=11, color=SUBTEXT, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════
# Slide 2: 페이지 구조
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
header(slide, "페이지 구조", "Next.js App Router 기반 페이지 목록")

pages = [
    ("/",                    "홈페이지",       "AI 에이전트 서비스 소개 및 사용법 안내",           ACCENT),
    ("/login",               "로그인",         "카카오 OAuth 소셜 로그인",                        ACCENT),
    ("/account/signup",      "회원가입",       "약관 동의 → 계정 생성",                           ACCENT),
    ("/auth-callback",       "OAuth 콜백",     "카카오 인증 후 토큰 저장 및 리다이렉트",           ACCENT),
    ("/stock-recommendation","주식 분석",      "종목코드 + 질문 입력 → 3개 에이전트 분석 결과",   GREEN),
    ("/board",               "게시판 목록",    "페이지네이션 게시물 목록",                         YELLOW),
    ("/board/create",        "게시물 작성",    "새 게시물 작성 폼",                               YELLOW),
    ("/board/read/[id]",     "게시물 상세",    "게시물 조회 / 수정 / 삭제",                       YELLOW),
    ("/youtube",             "YouTube 피드",   "주식 관련 유튜브 영상 목록",                       PURPLE),
]

for i, (path, name, desc, col) in enumerate(pages):
    row = i % 5
    col_idx = i // 5
    lx = 0.35 + col_idx * 6.55
    ty = 1.3 + row * 1.18

    rc(slide, lx, ty, 6.2, 1.05, CARD_BG)
    rc(slide, lx, ty, 0.06, 1.05, col)
    tb(slide, path, lx + 0.15, ty + 0.05, 2.2, 0.38, size=10, bold=True, color=col)
    tb(slide, name, lx + 2.4, ty + 0.05, 1.5, 0.38, size=11, bold=True, color=DARK_TEXT)
    tb(slide, desc, lx + 0.15, ty + 0.55, 5.9, 0.38, size=10, color=SUBTEXT)


# ══════════════════════════════════════════════
# Slide 3: 주식 분석 페이지 (핵심 기능)
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
header(slide, "주식 분석 페이지 (핵심 기능)", "/stock-recommendation — 메인 AI 분석 화면")

# 입력 섹션
rc(slide, 0.35, 1.3, 4.0, 5.7, CARD_BG)
rc(slide, 0.35, 1.3, 4.0, 0.42, ACCENT)
tb(slide, "입력 & 흐름", 0.35, 1.3, 4.0, 0.42,
   size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

flow_steps = [
    ("①", "종목 코드 입력", "예: 005930 (삼성전자)"),
    ("②", "질문 입력", "예: 투자해도 될까요?"),
    ("③", "분석 요청", "POST /api/v1/agent/query"),
    ("④", "로딩 단계 표시", "3단계 진행 애니메이션 (7초 간격)"),
    ("⑤", "결과 렌더링", "종합 신호 + 에이전트 카드 3개"),
    ("⑥", "이력 조회", "GET /api/v1/agent/history"),
]
for i, (num, title, desc) in enumerate(flow_steps):
    ty = 1.88 + i * 0.82
    rc(slide, 0.45, ty, 3.75, 0.72, ROW_ALT if i % 2 else CARD_BG)
    tb(slide, num, 0.55, ty + 0.06, 0.35, 0.55, size=13, bold=True, color=ACCENT)
    tb(slide, title, 0.95, ty + 0.05, 2.5, 0.32, size=11, bold=True, color=DARK_TEXT)
    tb(slide, desc,  0.95, ty + 0.37, 3.0, 0.28, size=9,  color=SUBTEXT)

# 결과 섹션
rc(slide, 4.65, 1.3, 8.35, 5.7, CARD_BG)
rc(slide, 4.65, 1.3, 8.35, 0.42, GREEN)
tb(slide, "결과 UI 컴포넌트", 4.65, 1.3, 8.35, 0.42,
   size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

components = [
    ("AnalysisResultHeader",  "종합 신호 배지 (bullish/bearish/neutral) + 신뢰도",  GREEN),
    ("AgentCard × 3",         "뉴스 / 공시 / 재무 에이전트 결과 카드",               ACCENT),
    ("SignalBadge",           "에이전트별 시그널 색상 배지",                         GREEN),
    ("ConfidenceBar",         "신뢰도 0.0~1.0 시각화 바",                           ACCENT),
    ("FinanceDataTable",      "재무 에이전트: ROE·ROA·부채비율·매출 표",             YELLOW),
    ("DisclosureDetails",     "공시 에이전트: 핵심 공시 목록 + 카테고리별 건수",      YELLOW),
    ("HistoryTimeline",       "이전 분석 이력 타임라인",                             PURPLE),
    ("AnalysisLoadingSteps",  "3단계 로딩 진행 표시 (7초 간격 애니메이션)",           SUBTEXT),
]
for i, (name, desc, col) in enumerate(components):
    ty = 1.88 + i * 0.72
    bg = ROW_ALT if i % 2 else CARD_BG
    rc(slide, 4.75, ty, 8.1, 0.62, bg)
    tb(slide, name, 4.85, ty + 0.06, 2.8, 0.45, size=11, bold=True, color=col)
    tb(slide, desc, 7.7,  ty + 0.06, 5.0, 0.45, size=10, color=SUBTEXT)


# ══════════════════════════════════════════════
# Slide 4: 게시판 & 인증
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
header(slide, "게시판 & 인증 기능")

# 게시판
rc(slide, 0.35, 1.3, 6.0, 5.7, CARD_BG)
rc(slide, 0.35, 1.3, 6.0, 0.42, YELLOW)
tb(slide, "게시판", 0.35, 1.3, 6.0, 0.42,
   size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

board_items = [
    ("BoardList",        "페이지네이션 게시물 목록 (10개/페이지)\n빈 상태 처리 포함"),
    ("BoardListItem",    "게시물 아이템 (제목·작성자·날짜)"),
    ("BoardCreateForm",  "게시물 작성 폼 (제목 + 본문)"),
    ("BoardEditForm",    "게시물 수정 폼"),
    ("BoardPostDetail",  "게시물 상세 + 수정/삭제 버튼"),
]
for i, (name, desc) in enumerate(board_items):
    ty = 1.9 + i * 0.98
    bg = ROW_ALT if i % 2 else CARD_BG
    rc(slide, 0.45, ty, 5.75, 0.85, bg)
    tb(slide, name, 0.55, ty + 0.05, 5.5, 0.35, size=12, bold=True, color=YELLOW)
    tb(slide, desc, 0.55, ty + 0.42, 5.5, 0.38, size=10, color=SUBTEXT)

# 인증
rc(slide, 6.85, 1.3, 6.1, 5.7, CARD_BG)
rc(slide, 6.85, 1.3, 6.1, 0.42, ACCENT)
tb(slide, "인증 (Auth)", 6.85, 1.3, 6.1, 0.42,
   size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

auth_items = [
    ("KakaoLoginButton",    "카카오 OAuth 로그인 버튼\n→ /kakao-authentication/request-oauth-link"),
    ("SignupContent",       "회원가입 폼\n(닉네임·이메일·비밀번호)"),
    ("TermsContent",        "약관 동의 UI\n(서비스·개인정보·아동보호)"),
    ("useAuthCallback",     "OAuth 콜백 훅\n토큰 저장 → 홈 리다이렉트"),
    ("useAuth / authAtom",  "Jotai 전역 인증 상태\nUNAUTHENTICATED / AUTHENTICATED"),
]
for i, (name, desc) in enumerate(auth_items):
    ty = 1.9 + i * 0.98
    bg = ROW_ALT if i % 2 else CARD_BG
    rc(slide, 6.95, ty, 5.85, 0.85, bg)
    tb(slide, name, 7.05, ty + 0.05, 5.6, 0.35, size=12, bold=True, color=ACCENT)
    tb(slide, desc, 7.05, ty + 0.42, 5.6, 0.38, size=10, color=SUBTEXT)


# ══════════════════════════════════════════════
# Slide 5: 아키텍처 & 상태 관리
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_bg(slide)
header(slide, "프론트엔드 아키텍처 & 상태 관리",
       "Clean Architecture + Jotai 원자적 상태 관리")

# Clean Architecture
rc(slide, 0.35, 1.3, 6.0, 5.7, CARD_BG)
rc(slide, 0.35, 1.3, 6.0, 0.42, PURPLE)
tb(slide, "Clean Architecture (features/)", 0.35, 1.3, 6.0, 0.42,
   size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

layers = [
    ("ui/components/",    "React 컴포넌트 (View)",                  ACCENT),
    ("application/",      "hooks + atoms (비즈니스 로직)",           GREEN),
    ("domain/",           "타입 · 모델 · 상태 정의",                YELLOW),
    ("infrastructure/",   "API 호출 · 로컬 스토리지",               PURPLE),
]
for i, (path, desc, col) in enumerate(layers):
    ty = 1.92 + i * 0.88
    rc(slide, 0.45, ty, 5.75, 0.75, ROW_ALT if i % 2 else CARD_BG)
    rc(slide, 0.45, ty, 0.06, 0.75, col)
    tb(slide, path, 0.6, ty + 0.06, 2.5, 0.3, size=12, bold=True, color=col)
    tb(slide, desc, 0.6, ty + 0.38, 5.5, 0.3, size=10, color=SUBTEXT)

tb(slide, "기능별 독립 모듈: auth / board / stock-recommendation / youtube",
   0.45, 5.62, 5.75, 0.45, size=10, color=SUBTEXT)

# Jotai 상태 관리
rc(slide, 6.85, 1.3, 6.1, 5.7, CARD_BG)
rc(slide, 6.85, 1.3, 6.1, 0.42, TEAL)
tb(slide, "Jotai 상태 관리", 6.85, 1.3, 6.1, 0.42,
   size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

atoms = [
    ("authAtom",            "인증 상태\nUNAUTHENTICATED | AUTHENTICATED",              ACCENT),
    ("stockAnalysisAtom",   "주식 분석 결과\nIDLE → ANALYZING → SUCCESS/ERROR",         GREEN),
    ("analysisHistoryAtom", "분석 이력 목록",                                          GREEN),
    ("boardAtom",           "게시판 목록 + 로딩/에러 상태",                            YELLOW),
]
for i, (name, desc, col) in enumerate(atoms):
    ty = 1.92 + i * 0.95
    bg = ROW_ALT if i % 2 else CARD_BG
    rc(slide, 6.95, ty, 5.85, 0.82, bg)
    tb(slide, name, 7.05, ty + 0.05, 5.6, 0.32, size=12, bold=True, color=col)
    tb(slide, desc, 7.05, ty + 0.4,  5.6, 0.35, size=10, color=SUBTEXT)

tb(slide, "HTTP Client (agentHttpClient.ts)",
   6.95, 5.72, 5.85, 0.32, size=11, bold=True, color=TEAL)
tb(slide, "Bearer 토큰 자동 주입 · 401 시 /login 자동 리다이렉트",
   6.95, 6.05, 5.85, 0.32, size=10, color=SUBTEXT)


# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
output = "docs/frontend-slides.pptx"
prs.save(output)
print(f"Saved: {output}")
